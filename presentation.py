# Create a professional PowerPoint presentation (15–20 slides) based on the project report

from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Edge AI-Based Grape Leaf Disease Detection"
slide.placeholders[1].text = (
    "ESP32-S3 Real-Time Vineyard Monitoring System\n"
    "Academic / Industry Project Presentation\nDecember 2025"
)

# Slides content
add_slide(
    "Project Motivation & Objective",
    "• Need for early disease detection in vineyards\n"
    "• Limitations of cloud-based monitoring (latency, cost, connectivity)\n"
    "• Objective: Develop a low-cost, real-time edge AI system\n"
    "• Fully autonomous operation without cloud dependency"
)

add_slide(
    "System Overview",
    "• Edge AI device deployed directly in vineyard\n"
    "• ESP32-S3 microcontroller with camera\n"
    "• On-device deep learning inference\n"
    "• Periodic monitoring with ultra-low power consumption"
)

add_slide(
    "Hardware Platform",
    "• ESP32-S3 (Dual-core LX7 @ 240 MHz)\n"
    "• 16 MB Flash, 8 MB PSRAM\n"
    "• OV3660 camera module (320×240, JPEG)\n"
    "• USB-C power and programming"
)

add_slide(
    "Software Stack",
    "• ESP-IDF v5.3.3\n"
    "• ESP-DL deep learning framework\n"
    "• Arduino ESP32 Camera Library (hybrid approach)\n"
    "• FreeRTOS-based real-time execution"
)

add_slide(
    "Development Methodology",
    "• Phase 1: Isolated component testing\n"
    "• Phase 2: Iterative system integration\n"
    "• Phase 3: Hybrid architecture and optimization\n"
    "• Continuous profiling and debugging"
)

add_slide(
    "AI Model Architecture",
    "• ESPDet-Pico lightweight object detector\n"
    "• INT8 quantized model (491 KB)\n"
    "• Input: 224×224 RGB images\n"
    "• Output: Bounding boxes and confidence scores"
)

add_slide(
    "Inference Pipeline",
    "1. Image capture (JPEG)\n"
    "2. JPEG decode to RGB888\n"
    "3. Resize/crop to model input\n"
    "4. AI inference on ESP32-S3\n"
    "5. Post-processing and detection output"
)

add_slide(
    "Memory Management Strategy",
    "• Model weights stored in DRAM\n"
    "• Camera buffers allocated in PSRAM\n"
    "• JPEG compression reduces memory usage by 94%\n"
    "• <10% PSRAM utilization during runtime"
)

add_slide(
    "Key Technical Challenges",
    "• Cache coherency conflicts (PSRAM vs Flash)\n"
    "• Limited DRAM availability\n"
    "• Camera initialization instability\n"
    "• Real-time performance constraints"
)

add_slide(
    "Key Innovations",
    "• Load AI model before camera initialization\n"
    "• Hybrid Arduino + ESP-IDF architecture\n"
    "• JPEG-based capture pipeline\n"
    "• Optimized memory hierarchy usage"
)

add_slide(
    "Performance Results",
    "• Inference time: 135–138 ms\n"
    "• End-to-end pipeline: ~182 ms\n"
    "• Frame rate: 5.3–5.5 FPS\n"
    "• Stable real-time operation"
)

add_slide(
    "Detection Accuracy",
    "• 5–7 detections per frame\n"
    "• Confidence range: 24–59%\n"
    "• Conservative detection behavior\n"
    "• Low observed false positives"
)

add_slide(
    "Power Efficiency",
    "• Active power: ~300 mW\n"
    "• Sleep power: ~10 mW\n"
    "• Average power: ~20 mW\n"
    "• >500 hours operation on 10 Ah battery"
)

add_slide(
    "System Stability",
    "• Zero crashes during testing\n"
    "• No memory leaks detected\n"
    "• Stable thermal behavior\n"
    "• Suitable for long-term deployment"
)

add_slide(
    "Applications & Use Cases",
    "• Early grape leaf disease detection\n"
    "• Vineyard health monitoring\n"
    "• Solar-powered edge deployments\n"
    "• Low-cost scalable sensor networks"
)

add_slide(
    "Limitations & Future Work",
    "• Improve detection confidence via retraining\n"
    "• Add data logging and wireless connectivity\n"
    "• Disease classification and severity estimation\n"
    "• Multi-sensor environmental integration"
)

add_slide(
    "Conclusions",
    "• Demonstrated feasibility of edge AI on ESP32-S3\n"
    "• Achieved real-time, low-power disease detection\n"
    "• Cost-effective and scalable solution\n"
    "• Strong foundation for commercial deployment"
)

# Save presentation
file_path = "/home/ubuntu/edge-ai-vineyard-monitoring/Edge_AI_Grape_Leaf_Disease_Detection_Presentation.pptx"
prs.save(file_path)

print(f"✅ Presentation created successfully!")
print(f"📁 Saved to: {file_path}")
